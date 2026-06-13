from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parent
ASSET_DIR = PROJECT_DIR / "report_assets"
OUT = PROJECT_DIR / "NLP_Mini_Project_1_Canva_Style_Final_Presentation.pptx"

NAVY = RGBColor(18, 39, 73)
BLUE = RGBColor(25, 96, 190)
SKY = RGBColor(219, 236, 255)
TEAL = RGBColor(20, 155, 150)
GREEN = RGBColor(43, 158, 108)
ORANGE = RGBColor(238, 142, 45)
RED = RGBColor(203, 69, 85)
INK = RGBColor(36, 42, 55)
GRAY = RGBColor(91, 101, 118)
LIGHT = RGBColor(247, 250, 255)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(212, 224, 240)


def font(run, size=18, bold=False, color=INK, name="Calibri"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=INK, align=None, name="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size=size, bold=bold, color=color, name=name)
    return box


def add_bullets(slide, items, x, y, w, h, size=17, color=INK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(6)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
    return box


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.18), Inches(2.6), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def title(slide, text, subtitle=None, num=None):
    add_bg(slide)
    add_text(slide, text, Inches(0.55), Inches(0.42), Inches(10.9), Inches(0.5), 26, True, NAVY)
    if subtitle:
        add_text(slide, subtitle, Inches(0.58), Inches(0.95), Inches(11.5), Inches(0.38), 12.5, False, GRAY)
    if num:
        add_text(slide, f"{num:02d}", Inches(12.45), Inches(7.06), Inches(0.5), Inches(0.24), 10, False, GRAY, PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, heading, body=None, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_text(slide, heading, x + Inches(0.2), y + Inches(0.15), w - Inches(0.35), Inches(0.3), 15, True, NAVY)
    if body:
        add_text(slide, body, x + Inches(0.2), y + Inches(0.52), w - Inches(0.35), h - Inches(0.62), 12.5, False, GRAY)
    return shape


def metric(slide, x, y, w, h, value, label, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    add_text(slide, value, x + Inches(0.15), y + Inches(0.18), w - Inches(0.3), Inches(0.4), 25, True, color, PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(0.12), y + Inches(0.65), w - Inches(0.24), Inches(0.34), 10.5, False, GRAY, PP_ALIGN.CENTER)


def image(slide, path, x, y, w=None, h=None):
    if path.exists():
        if w and h:
            slide.shapes.add_picture(str(path), x, y, width=w, height=h)
        elif w:
            slide.shapes.add_picture(str(path), x, y, width=w)
        else:
            slide.shapes.add_picture(str(path), x, y, height=h)


def table(slide, rows, cols, x, y, w, h, data, header_fill=BLUE):
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = Pt(5)
            cell.margin_right = Pt(5)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10.5)
                    p.font.color.rgb = INK
    return tbl


def connector(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = BLUE
    line.line.width = Pt(2)
    return line


def temple_illustration(slide, x, y, scale=1.0):
    base_w = Inches(4.0) * scale
    base_h = Inches(0.35) * scale
    colors = [NAVY, BLUE, TEAL]
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(2.25) * scale, base_w, base_h).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = NAVY
    slide.shapes[-1].line.fill.background()
    for i, (cx, h, col) in enumerate([
        (0.28, 1.25, BLUE), (1.05, 1.65, TEAL), (1.85, 2.15, NAVY), (2.65, 1.65, TEAL), (3.42, 1.25, BLUE)
    ]):
        left = x + Inches(cx) * scale
        top = y + Inches(2.25 - h) * scale
        tower = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, Inches(0.5) * scale, Inches(h) * scale)
        tower.fill.solid()
        tower.fill.fore_color.rgb = col
        tower.line.fill.background()
    for j in range(5):
        step = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.25 + j * 0.7) * scale, y + Inches(2.0) * scale, Inches(0.35) * scale, Inches(0.28) * scale)
        step.fill.solid()
        step.fill.fore_color.rgb = SKY
        step.line.fill.background()


def token_box(slide, text, x, y, color=WHITE, border=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.55), Inches(0.48))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = border
    add_text(slide, text, x + Inches(0.05), y + Inches(0.12), Inches(1.45), Inches(0.22), 13, True, NAVY, PP_ALIGN.CENTER, "Khmer OS Siemreap")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_text(s, "NLP Mini Project 1", Inches(0.7), Inches(0.75), Inches(5.7), Inches(0.45), 20, True, BLUE)
    add_text(s, "Khmer 4-Gram Text Generation", Inches(0.68), Inches(1.18), Inches(7.1), Inches(0.75), 34, True, NAVY)
    add_text(s, "Backoff vs Interpolation Language Models using Angkor-related Khmer Wikipedia corpus", Inches(0.72), Inches(2.02), Inches(6.9), Inches(0.55), 16, False, GRAY)
    temple_illustration(s, Inches(8.1), Inches(1.1), 1.05)
    metric(s, Inches(0.72), Inches(5.2), Inches(2.15), Inches(1.05), "19,681", "tokens after preprocessing", BLUE)
    metric(s, Inches(3.05), Inches(5.2), Inches(2.15), Inches(1.05), "500", "vocabulary cap", TEAL)
    metric(s, Inches(5.38), Inches(5.2), Inches(2.15), Inches(1.05), "27.22", "best test perplexity", GREEN)
    add_text(s, "01", Inches(12.45), Inches(7.06), Inches(0.5), Inches(0.24), 10, False, GRAY, PP_ALIGN.RIGHT)

    # 2
    s = prs.slides.add_slide(blank)
    title(s, "Presentation Roadmap", "A 20-minute walkthrough from data to model results.", 2)
    cards = [
        ("1. Corpus", "Angkor Wat and related Khmer articles", BLUE),
        ("2. Preprocess", "Clean text and tokenize Khmer words", TEAL),
        ("3. Models", "LM1 Backoff and LM2 Interpolation", ORANGE),
        ("4. Results", "Perplexity, vocabulary trade-off, demo", GREEN),
    ]
    x = Inches(0.8)
    for i, (h, b, c) in enumerate(cards):
        card(s, x + Inches(i * 3.05), Inches(2.0), Inches(2.65), Inches(1.5), h, b, c)
        if i < 3:
            connector(s, x + Inches(i * 3.05 + 2.65), Inches(2.75), x + Inches((i + 1) * 3.05), Inches(2.75))
    add_bullets(s, [
        "The deck follows the assignment structure, then explains the experimental result.",
        "Graphs and tables are placed near the section they support.",
        "The final slides explain why LM1 performs better and what can be improved."
    ], Inches(1.0), Inches(4.4), Inches(10.9), Inches(1.3), 18)

    # 3
    s = prs.slides.add_slide(blank)
    title(s, "Assignment Requirement Mapping", "How the notebook satisfies the NLP mini project instructions.", 3)
    data = [
        ["Requirement", "Implementation in this project"],
        ["Corpus", "Five Khmer Wikipedia articles about Angkor and Khmer history"],
        ["Split", "70% train, 10% validation, 20% test"],
        ["Tokenization", "khmernltk.word_tokenize with fallback option"],
        ["Vocabulary", "Top 500 words kept; rare words mapped to <UNK>"],
        ["Models", "LM1 Backoff 4-gram and LM2 Interpolation 4-gram"],
        ["Evaluation", "Perplexity on validation/test data"],
        ["Generator", "Interactive seed-based text generation demo"],
    ]
    table(s, len(data), 2, Inches(0.85), Inches(1.6), Inches(11.6), Inches(4.8), data)

    # 4
    s = prs.slides.add_slide(blank)
    title(s, "Corpus Source and Scope", "A focused historical corpus improves repeated phrase patterns.", 4)
    add_bullets(s, [
        "Main topic: Angkor Wat and related Khmer historical/cultural articles.",
        "Articles used: Angkor Wat, Angkor Thom, Bayon Temple, Siem Reap Province, Khmer Empire.",
        "Raw fallback corpus size: about 96,199 characters.",
        "The fallback corpus lets the notebook run even when scraping is unavailable."
    ], Inches(0.75), Inches(1.6), Inches(6.0), Inches(2.5), 18)
    temple_illustration(s, Inches(8.2), Inches(1.15), 0.95)
    card(s, Inches(7.35), Inches(4.65), Inches(4.95), Inches(1.25), "Why this corpus?", "The articles share words such as ប្រាសាទ, អង្គរ, ព្រះបាទ, ឆ្នាំ, and កម្ពុជា, giving the n-gram models stronger local patterns.", TEAL)

    # 5
    s = prs.slides.add_slide(blank)
    title(s, "Preprocessing Pipeline", "Turning raw article text into tokens for language modeling.", 5)
    steps = [
        ("Raw Article", "Wikipedia or fallback corpus"),
        ("Clean Text", "Remove labels and English-only lines"),
        ("Tokenize", "Khmer word segmentation"),
        ("Vocabulary", "Keep top 500 tokens"),
        ("Model Data", "Replace rare words with <UNK>"),
    ]
    x0 = Inches(0.7)
    for i, (h, b) in enumerate(steps):
        card(s, x0 + Inches(i * 2.5), Inches(2.0), Inches(2.05), Inches(1.35), h, b, [BLUE, TEAL, GREEN, ORANGE, RED][i])
        if i < 4:
            connector(s, x0 + Inches(i * 2.5 + 2.05), Inches(2.67), x0 + Inches((i + 1) * 2.5), Inches(2.67))
    add_text(s, "Key preprocessing choices", Inches(0.8), Inches(4.2), Inches(4.8), Inches(0.35), 18, True, NAVY)
    add_bullets(s, [
        "Keep punctuation separation so sentence endings can be learned.",
        "Do not remove stop words because language modeling needs grammar words.",
        "Use Khmer segmentation to improve token quality compared with plain split."
    ], Inches(0.85), Inches(4.75), Inches(11.4), Inches(1.25), 17)

    # 6
    s = prs.slides.add_slide(blank)
    title(s, "Dataset Split Result", "The 70/10/20 split follows the assignment requirement.", 6)
    image(s, ASSET_DIR / "human_dataset_split.png", Inches(0.8), Inches(1.55), Inches(6.1), Inches(4.5))
    metric(s, Inches(7.35), Inches(1.65), Inches(2.25), Inches(1.0), "13,776", "training tokens", BLUE)
    metric(s, Inches(9.95), Inches(1.65), Inches(2.25), Inches(1.0), "1,968", "validation tokens", ORANGE)
    metric(s, Inches(7.35), Inches(3.05), Inches(2.25), Inches(1.0), "3,937", "testing tokens", GREEN)
    metric(s, Inches(9.95), Inches(3.05), Inches(2.25), Inches(1.0), "19,681", "total tokens", TEAL)
    card(s, Inches(7.35), Inches(4.65), Inches(4.85), Inches(1.1), "Why validation matters", "Validation data is used to choose LM2 lambda and k values before evaluating on the test set.", BLUE)

    # 7
    s = prs.slides.add_slide(blank)
    title(s, "Vocabulary and <UNK>", "Vocabulary limitation is required by the assignment.", 7)
    add_bullets(s, [
        "Vocabulary cap: 500 most frequent tokens.",
        "Words outside the vocabulary become <UNK>.",
        "<UNK> means unknown token, not one exact recoverable word.",
        "Validation <UNK>: 660 tokens; Test <UNK>: 1,433 tokens."
    ], Inches(0.75), Inches(1.55), Inches(5.8), Inches(2.1), 18)
    data = [
        ["Token Type", "Example / Meaning"],
        ["Vocabulary word", "ប្រាសាទ, អង្គរវត្ត, មាន, ក្នុង"],
        ["Rare word", "Low-frequency names, numbers, long terms"],
        ["<UNK>", "Shared placeholder for rare/out-of-vocabulary words"],
    ]
    table(s, len(data), 2, Inches(6.95), Inches(1.55), Inches(5.55), Inches(2.25), data)
    add_text(s, "Generated example", Inches(0.78), Inches(4.25), Inches(4.5), Inches(0.3), 17, True, NAVY)
    token_box(s, "ប្រាសាទ", Inches(0.85), Inches(4.85), SKY)
    token_box(s, "អង្គរវត្ត", Inches(2.55), Inches(4.85), SKY)
    token_box(s, "ត្រូវបាន", Inches(4.25), Inches(4.85), SKY)
    token_box(s, "<UNK>", Inches(5.95), Inches(4.85), RGBColor(255, 232, 235), RED)
    token_box(s, "ក្នុង", Inches(7.65), Inches(4.85), SKY)
    token_box(s, "ឆ្នាំ", Inches(9.35), Inches(4.85), SKY)
    token_box(s, "<UNK>", Inches(11.05), Inches(4.85), RGBColor(255, 232, 235), RED)

    # 8
    s = prs.slides.add_slide(blank)
    title(s, "How 4-Gram Generation Works", "Both models use previous words to predict the next token.", 8)
    add_text(s, "A 4-gram sequence has four tokens:", Inches(0.85), Inches(1.5), Inches(5.2), Inches(0.32), 18, True, NAVY)
    token_box(s, "ប្រាសាទ", Inches(1.0), Inches(2.15), SKY)
    token_box(s, "អង្គរវត្ត", Inches(2.8), Inches(2.15), SKY)
    token_box(s, "មាន", Inches(4.6), Inches(2.15), SKY)
    token_box(s, "កម្ពស់", Inches(6.4), Inches(2.15), RGBColor(226, 248, 238), GREEN)
    add_text(s, "previous 3 words", Inches(1.0), Inches(2.85), Inches(5.1), Inches(0.28), 13, False, GRAY, PP_ALIGN.CENTER)
    add_text(s, "predicted next word", Inches(6.2), Inches(2.85), Inches(2.0), Inches(0.28), 13, False, GREEN, PP_ALIGN.CENTER)
    add_bullets(s, [
        "The seed provides the starting context.",
        "The generator predicts one token, appends it, then repeats.",
        "If the seed has fewer than three words, shorter contexts are used at the beginning."
    ], Inches(0.9), Inches(4.1), Inches(11.4), Inches(1.35), 18)

    # 9
    s = prs.slides.add_slide(blank)
    title(s, "LM1: Unsmoothed Backoff 4-Gram", "Use the most specific context first, then back off if needed.", 9)
    levels = [("4-gram", "w1 w2 w3 -> next"), ("3-gram", "w2 w3 -> next"), ("2-gram", "w3 -> next"), ("1-gram", "overall frequency")]
    for i, (h, b) in enumerate(levels):
        card(s, Inches(1.0 + i * 2.9), Inches(2.0 + i * 0.35), Inches(2.35), Inches(1.0), h, b, [BLUE, TEAL, ORANGE, GREEN][i])
    add_bullets(s, [
        "LM1 is strong when the corpus contains repeated phrases.",
        "It performed best in this experiment with test perplexity 27.2196.",
        "Lower perplexity means the model predicts the test text better."
    ], Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.25), 18)

    # 10
    s = prs.slides.add_slide(blank)
    title(s, "LM2: Interpolation with Add-k Smoothing", "Combine several probability levels instead of choosing only one.", 10)
    add_text(s, "P(next) = λ1·P1 + λ2·P2 + λ3·P3 + λ4·P4", Inches(0.9), Inches(1.6), Inches(7.2), Inches(0.45), 24, True, BLUE)
    data = [
        ["Parameter", "Selected value", "Meaning"],
        ["lambdas", "(0.4, 0.3, 0.2, 0.1)", "40% unigram, 30% bigram, 20% trigram, 10% 4-gram"],
        ["k", "0.001", "Small add-k smoothing value"],
        ["Validation PP", "47.4493", "Best validation score found in grid search"],
        ["Test PP", "45.7879", "Final performance on unseen test data"],
    ]
    table(s, len(data), 3, Inches(0.9), Inches(2.45), Inches(11.6), Inches(2.8), data)
    card(s, Inches(1.1), Inches(5.75), Inches(10.9), Inches(0.75), "Interpretation", "LM2 is theoretically robust, but in this corpus smoothing and sparsity made it less accurate than LM1.", ORANGE)

    # 11
    s = prs.slides.add_slide(blank)
    title(s, "Hyperparameter Tuning", "LM2 parameters were selected using validation perplexity.", 11)
    add_bullets(s, [
        "Tested multiple lambda combinations for unigram, bigram, trigram, and 4-gram weights.",
        "Tested k values: 0.001, 0.005, 0.01, 0.02, 0.05, 0.10.",
        "Selected the setting with the lowest validation perplexity.",
        "Best result: lambdas=(0.4, 0.3, 0.2, 0.1), k=0.001."
    ], Inches(0.85), Inches(1.55), Inches(6.2), Inches(2.5), 18)
    card(s, Inches(7.55), Inches(1.7), Inches(4.3), Inches(1.05), "Why tune?", "Different corpora need different smoothing and probability weights.", BLUE)
    card(s, Inches(7.55), Inches(3.1), Inches(4.3), Inches(1.05), "Selection rule", "Choose lowest validation perplexity before testing.", TEAL)
    card(s, Inches(7.55), Inches(4.5), Inches(4.3), Inches(1.05), "Final LM2", "Validation PP: 47.4493", GREEN)

    # 12
    s = prs.slides.add_slide(blank)
    title(s, "Main Result Table", "The final comparison shows LM1 is the better model for this corpus.", 12)
    data = [
        ["Model", "Method", "Best parameters", "Validation PP", "Test PP"],
        ["LM1", "Unsmoothed backoff 4-gram", "No lambda/k", "-", "27.2196"],
        ["LM2", "Interpolation + add-k", "λ=(0.4,0.3,0.2,0.1), k=0.001", "47.4493", "45.7879"],
    ]
    table(s, len(data), 5, Inches(0.55), Inches(1.8), Inches(12.25), Inches(1.6), data)
    metric(s, Inches(1.0), Inches(4.15), Inches(2.5), Inches(1.15), "27.2196", "LM1 test perplexity", GREEN)
    metric(s, Inches(4.05), Inches(4.15), Inches(2.5), Inches(1.15), "45.7879", "LM2 test perplexity", ORANGE)
    card(s, Inches(7.35), Inches(4.05), Inches(4.9), Inches(1.35), "Conclusion from table", "Lower perplexity is better, so LM1 Backoff is selected as the stronger final model for this experiment.", GREEN)

    # 13
    s = prs.slides.add_slide(blank)
    title(s, "Model Perplexity Graph", "Visual comparison of final test perplexity.", 13)
    image(s, ASSET_DIR / "human_model_perplexity.png", Inches(0.85), Inches(1.45), Inches(6.3), Inches(4.7))
    add_bullets(s, [
        "LM1 has lower perplexity than LM2.",
        "This means LM1 gives higher probability to the actual test sequence.",
        "The result does not mean interpolation is bad; it means this corpus/setting favored backoff."
    ], Inches(7.55), Inches(1.65), Inches(4.8), Inches(2.0), 18)
    card(s, Inches(7.55), Inches(4.4), Inches(4.7), Inches(1.15), "Key message", "For the final notebook result, LM1 is easier to defend because it has the best test perplexity.", BLUE)

    # 14
    s = prs.slides.add_slide(blank)
    title(s, "Vocabulary Size vs Perplexity", "Readability and probability are in tension.", 14)
    image(s, ASSET_DIR / "human_vocab_perplexity.png", Inches(0.65), Inches(1.35), Inches(7.0), Inches(4.95))
    add_bullets(s, [
        "Vocabulary sizes tested: 100, 200, 300, 500, 700, 1000.",
        "Larger vocabulary reduces <UNK> but creates sparser n-gram counts.",
        "In this experiment, perplexity increased when vocabulary became too large.",
        "Vocabulary cap 500 was chosen as a practical balance."
    ], Inches(7.95), Inches(1.55), Inches(4.7), Inches(2.7), 17)

    # 15
    s = prs.slides.add_slide(blank)
    title(s, "Text Generation Demo", "The demo lets a user test the model with a seed phrase.", 15)
    card(s, Inches(0.85), Inches(1.55), Inches(5.55), Inches(1.0), "Seed", "Example: ប្រាសាទ អង្គរវត្ត", BLUE)
    card(s, Inches(0.85), Inches(2.8), Inches(5.55), Inches(1.0), "Model Selector", "Choose LM1 Backoff or LM2 Interpolation", TEAL)
    card(s, Inches(0.85), Inches(4.05), Inches(5.55), Inches(1.0), "Length", "Controls how many tokens the model generates", ORANGE)
    card(s, Inches(7.0), Inches(1.55), Inches(5.2), Inches(1.25), "Raw Text", "Shows the real output, including <UNK>.", RED)
    card(s, Inches(7.0), Inches(3.15), Inches(5.2), Inches(1.25), "Clean Text", "Removes <UNK> only for easier reading.", GREEN)
    add_text(s, "Important: clean text is for display only; evaluation uses the real processed tokens.", Inches(7.0), Inches(5.15), Inches(5.3), Inches(0.6), 16, True, NAVY)

    # 16
    s = prs.slides.add_slide(blank)
    title(s, "Interpreting <UNK> in Generated Text", "The exact original word behind <UNK> cannot be recovered.", 16)
    add_text(s, "Generated:", Inches(0.9), Inches(1.55), Inches(2.0), Inches(0.28), 16, True, NAVY)
    add_text(s, "ប្រាសាទ អង្គរវត្ត ត្រូវបាន <UNK> ក្នុង ឆ្នាំ <UNK> ...", Inches(0.9), Inches(1.95), Inches(11.5), Inches(0.42), 19, False, INK, name="Khmer OS Siemreap")
    add_text(s, "Possible original article phrase:", Inches(0.9), Inches(2.85), Inches(4.0), Inches(0.28), 16, True, NAVY)
    add_text(s, "ប្រាសាទ អង្គរវត្ត ត្រូវបាន កសាងឡើង ក្នុង ឆ្នាំ ១១២២ ...", Inches(0.9), Inches(3.25), Inches(11.5), Inches(0.42), 19, False, INK, name="Khmer OS Siemreap")
    add_bullets(s, [
        "The possible phrase is inferred from article context.",
        "The model itself only predicted the <UNK> category.",
        "This limitation is caused by replacing many rare words with one shared token."
    ], Inches(0.95), Inches(4.45), Inches(11.2), Inches(1.25), 18)

    # 17
    s = prs.slides.add_slide(blank)
    title(s, "Why LM1 Performed Better", "The result is explainable from corpus size and sparsity.", 17)
    card(s, Inches(0.9), Inches(1.55), Inches(3.4), Inches(1.25), "Repeated phrases", "Backoff benefits when the article repeats exact local patterns.", BLUE)
    card(s, Inches(4.95), Inches(1.55), Inches(3.4), Inches(1.25), "Sparse 4-grams", "Khmer segmentation and rare words make many contexts infrequent.", ORANGE)
    card(s, Inches(9.0), Inches(1.55), Inches(3.4), Inches(1.25), "Smoothing spread", "LM2 gives small probability to many tokens, which can increase perplexity.", RED)
    add_bullets(s, [
        "LM2 is not automatically better in every corpus.",
        "Interpolation usually helps when lambda and smoothing match the data well.",
        "For this project, the validation-tuned LM2 still had higher test perplexity than LM1."
    ], Inches(1.05), Inches(4.0), Inches(11.1), Inches(1.45), 18)

    # 18
    s = prs.slides.add_slide(blank)
    title(s, "Limitations and Future Improvements", "The project is complete, but the model has natural n-gram limits.", 18)
    add_bullets(s, [
        "Generated Khmer text is sometimes not fully natural.",
        "<UNK> tokens reduce readability.",
        "Larger vocabulary reduces <UNK> but can worsen perplexity.",
        "N-gram models do not understand semantic meaning.",
        "Future improvements: larger corpus, stronger Khmer segmentation, better smoothing, and richer evaluation examples."
    ], Inches(0.85), Inches(1.55), Inches(6.1), Inches(3.2), 18)
    temple_illustration(s, Inches(8.1), Inches(2.2), 0.9)
    card(s, Inches(7.25), Inches(5.45), Inches(5.1), Inches(0.8), "Submission position", "The notebook stays strict to the assignment: 4-gram LM1, LM2, validation tuning, test perplexity, and generator.", GREEN)

    # 19
    s = prs.slides.add_slide(blank)
    title(s, "Conclusion", "The project successfully demonstrates classical Khmer n-gram language modeling.", 19)
    metric(s, Inches(0.9), Inches(1.65), Inches(2.35), Inches(1.05), "LM1", "best final model", GREEN)
    metric(s, Inches(3.55), Inches(1.65), Inches(2.35), Inches(1.05), "27.22", "lowest test perplexity", BLUE)
    metric(s, Inches(6.2), Inches(1.65), Inches(2.35), Inches(1.05), "500", "vocabulary cap", TEAL)
    metric(s, Inches(8.85), Inches(1.65), Inches(2.35), Inches(1.05), "4-gram", "previous 3 words -> next", ORANGE)
    add_bullets(s, [
        "Corpus quality, tokenization, vocabulary size, and smoothing strongly affect results.",
        "Backoff performed better than interpolation on this specific corpus.",
        "The interactive demo makes the model behavior clear through seed, length, raw text, and clean text."
    ], Inches(1.05), Inches(3.65), Inches(11.0), Inches(1.4), 19)
    add_text(s, "Thank you", Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.55), 28, True, NAVY, PP_ALIGN.CENTER)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(path)
