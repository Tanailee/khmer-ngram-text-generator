from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parent
ASSET_DIR = PROJECT_DIR / "report_assets"
OUT = PROJECT_DIR / "NLP_Mini_Project_1_20min_Presentation.pptx"

BLUE = RGBColor(23, 86, 159)
LIGHT_BLUE = RGBColor(226, 240, 253)
NAVY = RGBColor(20, 39, 78)
GREEN = RGBColor(38, 156, 111)
ORANGE = RGBColor(242, 151, 39)
RED = RGBColor(208, 76, 86)
GRAY = RGBColor(90, 98, 112)
WHITE = RGBColor(255, 255, 255)


def set_font(run, size=24, bold=False, color=NAVY):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, text, left, top, width, height, size=24, bold=False, color=NAVY, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.55), size=30, bold=True, color=NAVY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.96), Inches(1.4), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.55), Inches(1.08), Inches(12.2), Inches(0.35), size=13, color=GRAY)


def add_footer(slide, num):
    add_textbox(slide, f"{num:02d}", Inches(12.55), Inches(7.05), Inches(0.5), Inches(0.25), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, left, top, width, height, size=20, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_card(slide, title, body, left, top, width, height, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(210, 223, 238)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, title, left + Inches(0.22), top + Inches(0.16), width - Inches(0.35), Inches(0.3), size=16, bold=True)
    add_textbox(slide, body, left + Inches(0.22), top + Inches(0.55), width - Inches(0.35), height - Inches(0.62), size=12.5, color=GRAY)


def add_metric(slide, value, label, left, top, width, color=BLUE):
    add_textbox(slide, value, left, top, width, Inches(0.38), size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left, top + Inches(0.42), width, Inches(0.25), size=11.5, color=GRAY, align=PP_ALIGN.CENTER)


def add_flow(slide, labels, left, top, box_w, box_h, gap):
    for i, label in enumerate(labels):
        x = left + i * (box_w + gap)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else WHITE
        shp.line.color.rgb = BLUE
        add_textbox(slide, label, x + Inches(0.08), top + Inches(0.18), box_w - Inches(0.16), box_h - Inches(0.25), size=13.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_textbox(slide, "→", x + box_w + Inches(0.02), top + Inches(0.18), gap - Inches(0.04), box_h, size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)


def new_slide(prs, num, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 251, 255)
    add_title(slide, title, subtitle)
    add_footer(slide, num)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 249, 255)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    add_textbox(slide, "Text Generation with 4-gram Language Models", Inches(0.8), Inches(1.45), Inches(8.2), Inches(0.8), size=34, bold=True)
    add_textbox(slide, "NLP Mini Project 1 | Khmer Wikipedia Corpus: Angkor Wat Domain", Inches(0.85), Inches(2.32), Inches(8.4), Inches(0.35), size=16, color=GRAY)
    add_textbox(slide, "LM1 Backoff vs LM2 Interpolation + Add-k Smoothing", Inches(0.85), Inches(2.85), Inches(8.6), Inches(0.35), size=18, bold=True, color=BLUE)
    add_card(slide, "20-minute presentation", "Data → Preprocessing → 4-gram models → Tuning → Perplexity → Generation demo", Inches(0.85), Inches(4.2), Inches(6.4), Inches(1.1), accent=GREEN)
    add_metric(slide, "19,681", "tokens after Khmer segmentation", Inches(9.3), Inches(1.5), Inches(2.2), GREEN)
    add_metric(slide, "5", "related Wikipedia articles", Inches(9.3), Inches(2.65), Inches(2.2), BLUE)
    add_metric(slide, "27.22", "LM1 test perplexity", Inches(9.3), Inches(3.8), Inches(2.2), ORANGE)
    add_footer(slide, 1)

    # 2
    slide = new_slide(prs, 2, "Presentation Roadmap", "Designed for a 20-minute explanation")
    add_flow(slide, ["Problem", "Corpus", "Preprocess", "LM1", "LM2", "Results"], Inches(0.7), Inches(2.05), Inches(1.55), Inches(0.75), Inches(0.25))
    add_flow(slide, ["Vocab", "Generation", "Demo", "Limitations", "Conclusion"], Inches(1.65), Inches(4.05), Inches(1.65), Inches(0.75), Inches(0.35))
    add_bullets(slide, ["Suggested timing: 1-2 minutes per main concept", "Spend more time on model differences and result interpretation", "Use the demo near the end to make Part II concrete"], Inches(1.0), Inches(5.25), Inches(11.0), Inches(1.2), size=16)

    # 3
    slide = new_slide(prs, 3, "Project Requirement Mapping")
    rows = [
        ("Corpus", "5 Khmer Wikipedia articles"),
        ("Split", "70% training, 10% validation, 20% testing"),
        ("Tokenize", "khmernltk.word_tokenize + split fallback"),
        ("Vocabulary", "Cap = 500, remaining tokens → <UNK>"),
        ("Models", "LM1 Backoff, LM2 Interpolation + add-k"),
        ("Evaluation", "Test perplexity + generated examples"),
    ]
    for i, (a, b) in enumerate(rows):
        add_card(slide, a, b, Inches(0.8 + (i % 2) * 6.1), Inches(1.55 + (i // 2) * 1.25), Inches(5.5), Inches(0.95), accent=BLUE if i % 2 == 0 else GREEN)

    # 4
    slide = new_slide(prs, 4, "Corpus: Larger and Topic-Consistent")
    add_bullets(slide, ["ប្រាសាទ អង្គរវត្ត", "អង្គរធំ", "ប្រាសាទបាយ័ន", "ខេត្តសៀមរាប", "ចក្រភពខ្មែរ"], Inches(0.9), Inches(1.55), Inches(4.0), Inches(2.3), size=19)
    add_card(slide, "Why combine articles?", "One article was too small for a language model. Related articles increase n-gram evidence while keeping the generated text on the Angkor/Khmer history domain.", Inches(5.0), Inches(1.55), Inches(7.2), Inches(1.45), accent=GREEN)
    add_metric(slide, "96,199", "corpus characters", Inches(5.1), Inches(3.65), Inches(2.1), BLUE)
    add_metric(slide, "19,681", "tokens", Inches(7.5), Inches(3.65), Inches(2.1), GREEN)
    add_metric(slide, "13,776", "training tokens", Inches(9.9), Inches(3.65), Inches(2.1), ORANGE)

    # 5
    slide = new_slide(prs, 5, "Preprocessing Pipeline")
    add_flow(slide, ["Scrape", "Clean", "Segment", "Split", "Limit vocab", "Train"], Inches(0.55), Inches(1.7), Inches(1.65), Inches(0.75), Inches(0.2))
    add_card(slide, "Cleaning choices", "Removed English-only lines, article markers, scripts/references, and separated Khmer sentence punctuation.", Inches(0.9), Inches(3.1), Inches(5.6), Inches(1.25), accent=BLUE)
    add_card(slide, "Khmer tokenization", "Used khmernltk.word_tokenize because Khmer word boundaries are not always separated by spaces.", Inches(6.8), Inches(3.1), Inches(5.6), Inches(1.25), accent=GREEN)
    add_card(slide, "Assignment alignment", "Vocabulary is limited and rare/out-of-vocabulary tokens become <UNK>.", Inches(3.3), Inches(5.0), Inches(6.7), Inches(0.95), accent=ORANGE)

    # 6
    slide = new_slide(prs, 6, "What Does 4-gram Mean?")
    add_textbox(slide, "A 4-gram model predicts the next token using the previous 3 tokens.", Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.45), size=22, bold=True, color=BLUE)
    add_card(slide, "Example context", "ប្រាសាទ អង្គរវត្ត ជា ___", Inches(1.0), Inches(2.55), Inches(4.6), Inches(1.0), accent=BLUE)
    add_textbox(slide, "→", Inches(5.8), Inches(2.75), Inches(0.5), Inches(0.5), size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_card(slide, "Possible next token", "ប្រាសាទ", Inches(6.55), Inches(2.55), Inches(3.3), Inches(1.0), accent=GREEN)
    add_bullets(slide, ["The model learns from counted sequences in training data.", "It does not understand meaning; it estimates probability from observed patterns."], Inches(1.15), Inches(4.3), Inches(10.8), Inches(1.2), size=18)

    # 7
    slide = new_slide(prs, 7, "LM1: Backoff 4-gram")
    add_flow(slide, ["Try 4-gram", "If unseen", "Try 3-gram", "Try 2-gram", "Unigram"], Inches(0.8), Inches(2.0), Inches(1.75), Inches(0.8), Inches(0.35))
    add_card(slide, "Core idea", "Use the longest available context. Back off only when the higher-order n-gram was not observed.", Inches(0.9), Inches(3.55), Inches(5.7), Inches(1.35), accent=BLUE)
    add_card(slide, "Why it worked well", "With this corpus, exact local patterns are useful, and backoff avoids zero probability by moving to shorter histories.", Inches(6.9), Inches(3.55), Inches(5.7), Inches(1.35), accent=GREEN)

    # 8
    slide = new_slide(prs, 8, "LM2: Interpolation + Add-k")
    add_textbox(slide, "P = λ1P(unigram) + λ2P(bigram) + λ3P(trigram) + λ4P(4-gram)", Inches(0.9), Inches(1.55), Inches(11.7), Inches(0.5), size=21, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_card(slide, "Best lambdas", "(0.4, 0.3, 0.2, 0.1)\nMore weight on lower-order n-grams", Inches(1.0), Inches(2.55), Inches(3.7), Inches(1.4), accent=BLUE)
    add_card(slide, "Best k", "0.001\nSmall smoothing for sparse n-grams", Inches(4.95), Inches(2.55), Inches(3.7), Inches(1.4), accent=GREEN)
    add_card(slide, "Interpretation", "The validation set preferred stable lower-order probabilities because exact 4-grams are still sparse.", Inches(8.9), Inches(2.55), Inches(3.6), Inches(1.4), accent=ORANGE)

    # 9
    slide = new_slide(prs, 9, "Dataset Split Visualization")
    slide.shapes.add_picture(str(ASSET_DIR / "human_dataset_split.png"), Inches(1.1), Inches(1.55), width=Inches(7.1))
    add_card(slide, "Why this matters", "Training learns counts. Validation tunes λ and k. Testing gives final unbiased evaluation.", Inches(8.55), Inches(2.0), Inches(3.7), Inches(2.3), accent=BLUE)

    # 10
    slide = new_slide(prs, 10, "Final Perplexity Results")
    slide.shapes.add_picture(str(ASSET_DIR / "human_model_perplexity.png"), Inches(0.9), Inches(1.45), width=Inches(6.5))
    add_metric(slide, "27.2196", "LM1 Backoff", Inches(8.0), Inches(1.7), Inches(2.2), GREEN)
    add_metric(slide, "45.7879", "LM2 Interpolation", Inches(10.3), Inches(1.7), Inches(2.2), RED)
    add_card(slide, "Main finding", "Lower perplexity is better. LM1 performs better on this corpus, although LM2 is smoother in theory.", Inches(8.05), Inches(3.35), Inches(4.3), Inches(1.35), accent=GREEN)

    # 11
    slide = new_slide(prs, 11, "Vocabulary Trade-off")
    slide.shapes.add_picture(str(ASSET_DIR / "human_vocab_perplexity.png"), Inches(0.7), Inches(1.35), width=Inches(7.5))
    add_card(slide, "Small vocab", "Lower perplexity, but many <UNK> tokens.", Inches(8.55), Inches(1.65), Inches(3.8), Inches(1.05), accent=BLUE)
    add_card(slide, "Large vocab", "More readable words, but sparse counts and higher perplexity.", Inches(8.55), Inches(3.0), Inches(3.8), Inches(1.05), accent=ORANGE)
    add_card(slide, "Chosen balance", "Vocabulary size = 500", Inches(8.55), Inches(4.35), Inches(3.8), Inches(0.95), accent=GREEN)

    # 12
    slide = new_slide(prs, 12, "Text Generation Examples")
    add_card(slide, "Seed", "ប្រាសាទ", Inches(0.8), Inches(1.45), Inches(2.4), Inches(0.8), accent=BLUE)
    add_card(slide, "LM1 Backoff", "ប្រាសាទ នេះ កសាង នៅ ចុង សតវត្សរ៍ ទី ១២ និង ១៣។", Inches(0.8), Inches(2.65), Inches(5.8), Inches(1.35), accent=GREEN)
    add_card(slide, "LM2 Interpolation", "ប្រាសាទ នេះ បង្ហាញ ឲ្យដឹង ថា ត្រូវបាន ទៅជា <UNK> ...", Inches(6.9), Inches(2.65), Inches(5.8), Inches(1.35), accent=ORANGE)
    add_bullets(slide, ["LM1 is cleaner for demo.", "LM2 is valid but more mixed due to smoothing and interpolation.", "Keep length around 20-30 tokens."], Inches(1.05), Inches(4.75), Inches(10.6), Inches(1.2), size=18)

    # 13
    slide = new_slide(prs, 13, "Live Demo Plan")
    add_card(slide, "Recommended demo settings", "Seed: ប្រាសាទ អង្គរវត្ត\nModel: LM1 Backoff\nLength: 20-25", Inches(0.9), Inches(1.55), Inches(4.8), Inches(1.8), accent=GREEN)
    add_card(slide, "What to show", "Click Generate Text and explain that each new token is predicted from previous context.", Inches(6.1), Inches(1.55), Inches(5.8), Inches(1.8), accent=BLUE)
    add_flow(slide, ["Seed", "Context", "Next token", "Update context", "Repeat"], Inches(1.2), Inches(4.15), Inches(1.7), Inches(0.75), Inches(0.42))

    # 14
    slide = new_slide(prs, 14, "Why LM1 Beat LM2")
    add_card(slide, "Reason 1", "Backoff uses strong observed local patterns when available.", Inches(0.9), Inches(1.6), Inches(5.6), Inches(1.15), accent=GREEN)
    add_card(slide, "Reason 2", "Interpolation spreads probability across multiple n-gram levels.", Inches(6.8), Inches(1.6), Inches(5.6), Inches(1.15), accent=ORANGE)
    add_card(slide, "Reason 3", "Even with larger corpus, exact 4-grams remain sparse after Khmer segmentation.", Inches(0.9), Inches(3.2), Inches(5.6), Inches(1.15), accent=BLUE)
    add_card(slide, "Key message", "LM2 is not guaranteed to beat LM1. The best model depends on corpus, tokenization, vocabulary, and test data.", Inches(6.8), Inches(3.2), Inches(5.6), Inches(1.15), accent=RED)

    # 15
    slide = new_slide(prs, 15, "Limitations and Future Work")
    add_bullets(slide, ["Generated text is not perfectly natural because n-gram models do not understand meaning.", "Some <UNK> remains because vocabulary limiting is required.", "Readable output and low perplexity are a trade-off.", "Future work: larger corpus, better Khmer segmentation, or neural language models."], Inches(0.9), Inches(1.55), Inches(11.4), Inches(2.3), size=19)
    add_card(slide, "How to defend this", "The assignment is about implementing and evaluating 4-gram language models, not building a perfect modern chatbot.", Inches(2.2), Inches(5.0), Inches(8.8), Inches(1.0), accent=BLUE)

    # 16
    slide = new_slide(prs, 16, "Conclusion")
    add_bullets(slide, ["Built both required 4-gram models.", "Used validation to tune λ and k.", "Evaluated with perplexity on the test set.", "LM1 Backoff performed best: 27.2196 test perplexity.", "Created a working text generator demo."], Inches(1.0), Inches(1.55), Inches(10.8), Inches(2.8), size=20)
    add_textbox(slide, "Thank you", Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.6), size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
